
from pspy import so_dict
from pspipe_utils import  pspipe_list
import sys
from pptx import Presentation
from pptx.util import Inches
from operator import itemgetter


d = so_dict.so_dict()
d.read_from_file(sys.argv[1])


skip_EB = True

plot_dir = "plots/array_nulls"

spectra = ["TT", "TE", "TB", "ET", "BT", "EE", "EB", "BE", "BB"]

if skip_EB == True:
    tested_spectra = ["TT", "TE", "ET", "TB", "BT", "EE", "BB"]
    plot_dir = "plots/array_nulls_skip_EB"
else:
    tested_spectra = spectra
    plot_dir = "plots/array_nulls"


null_list = pspipe_list.get_null_list(d, spectra=tested_spectra)
null_list = sorted(null_list, key=itemgetter(0))[::-1]
n_null = len(null_list)
print(n_null)


prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
width = Inches(8)
height = Inches(6)

slide = prs.slides.add_slide(blank_slide_layout)
left = int((prs.slide_width-width)/2)
top = int((prs.slide_height-height)/2)

pic = slide.shapes.add_picture(f"{plot_dir}/pte_hist_all_corrected_spectra+mc+beam+leakage_cov.png", left, top,height=height,width=width)

for null in null_list:
    mode, ms1, ms2, ms3, ms4 = null
    fname = f"diff_{mode}_{ms1}x{ms2}_{ms3}x{ms4}"

    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(f"{plot_dir}/{fname}.png", left, top,height=height,width=width)

prs.save("null.pptx")
